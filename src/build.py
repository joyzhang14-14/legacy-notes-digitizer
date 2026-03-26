"""
Phase 4: 白纸合成 — Trace SVG + 数字文字 → PDF/PPT/纯文字
用法：
  python src/build.py                          # 生成全部输出
  python src/build.py --pages 1-3             # 只处理第1-3页
  python src/build.py --format pdf,pptx,text  # 指定输出格式
  python src/build.py --force                 # 强制覆盖已有输出
  python src/build.py --corrections corrections.json
"""

import argparse
import json
import math
import os
import re
import sys
import urllib.request
from pathlib import Path
from typing import Optional

# ─────────────────────────── 路径 ───────────────────────────

ROOT          = Path(__file__).parent.parent
PAGES_DIR     = ROOT / "intermediate" / "pages"
STRUCTURE_DIR = ROOT / "intermediate" / "structure"
OCR_DIR       = ROOT / "intermediate" / "text_ocr"
TRACES_DIR    = ROOT / "intermediate" / "traces"
OUTPUT_DIR    = ROOT / "output"
FONTS_DIR     = ROOT / "fonts"

# ─────────────────────────── 字体 ───────────────────────────

_NOTO_URLS = {
    "NotoSansSC-Regular": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTF/SimplifiedChinese/"
        "NotoSansCJKsc-Regular.otf"
    ),
    "NotoSansSC-Bold": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTF/SimplifiedChinese/"
        "NotoSansCJKsc-Bold.otf"
    ),
    "NotoSerifSC-Regular": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Serif/OTF/SimplifiedChinese/"
        "NotoSerifCJKsc-Regular.otf"
    ),
}

_MACOS_FALLBACKS = [
    ("STHeitiMedium", "/System/Library/Fonts/STHeiti Medium.ttc",  0, "sans-bold"),
    ("STHeitiLight",  "/System/Library/Fonts/STHeiti Light.ttc",   0, "sans"),
    ("Songti",        "/System/Library/Fonts/Supplemental/Songti.ttc", 1, "serif"),
]

_REGISTERED_FONTS: dict[str, str] = {}


def _try_register(rl_name: str, path: str, subfont_index: int = 0) -> bool:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    try:
        pdfmetrics.registerFont(TTFont(rl_name, path, subfontIndex=subfont_index))
        _REGISTERED_FONTS[rl_name] = path
        return True
    except Exception as e:
        print(f"  [字体] 注册 {rl_name} 失败: {e}")
        return False


def _download_noto() -> dict[str, str]:
    FONTS_DIR.mkdir(parents=True, exist_ok=True)
    downloaded = {}
    for name, url in _NOTO_URLS.items():
        path = FONTS_DIR / f"{name}.otf"
        if not path.exists():
            print(f"  [字体] 下载 {name}...")
            try:
                urllib.request.urlretrieve(url, str(path))
                print(f"  [字体] 下载完成: {name}")
            except Exception as e:
                print(f"  [字体] 下载失败 {name}: {e}")
        if path.exists():
            downloaded[name] = str(path)
    return downloaded


def setup_fonts(download_noto: bool = False) -> dict:
    print("[build] 初始化字体...")
    noto_available = {}
    for name in _NOTO_URLS:
        path = FONTS_DIR / f"{name}.otf"
        if path.exists() and _try_register(name, str(path)):
            noto_available[name] = str(path)

    if download_noto and len(noto_available) < len(_NOTO_URLS):
        dl = _download_noto()
        for name, path in dl.items():
            if name not in noto_available and _try_register(name, path):
                noto_available[name] = path

    if noto_available:
        import platform, shutil
        if platform.system() == "Darwin":
            user_fonts = Path.home() / "Library" / "Fonts"
            user_fonts.mkdir(parents=True, exist_ok=True)
            for name, src in noto_available.items():
                dst = user_fonts / Path(src).name
                if not dst.exists():
                    try:
                        shutil.copy2(src, dst)
                    except Exception:
                        pass

    system_fonts = {}
    for rl_name, path, idx, ftype in _MACOS_FALLBACKS:
        if os.path.exists(path) and _try_register(rl_name, path, idx):
            system_fonts[ftype] = rl_name

    def pick(noto_key: str, fallback_type: str) -> str:
        if noto_key in noto_available:
            return noto_key
        return system_fonts.get(fallback_type, system_fonts.get("sans", "Helvetica"))

    font_map = {
        "title":      (pick("NotoSansSC-Bold",     "sans-bold"), 28),
        "subtitle":   (pick("NotoSansSC-Bold",     "sans-bold"), 20),
        "body":       (pick("NotoSerifSC-Regular", "serif"),     14),
        "annotation": (pick("NotoSerifSC-Regular", "serif"),     12),
        "label":      (pick("NotoSansSC-Regular",  "sans"),      11),
        "dimension":  (pick("NotoSansSC-Regular",  "sans"),      10),
    }
    print(f"  [字体] title={font_map['title'][0]}, body={font_map['body'][0]}")
    return font_map


# ─────────────────────────── 颜色 ───────────────────────────

_COLOR_RGB = {
    "black": (0.0,  0.0,  0.0),
    "red":   (0.80, 0.08, 0.08),
    "blue":  (0.10, 0.20, 0.72),
}


def get_color_rgb(color_name: str) -> tuple:
    return _COLOR_RGB.get(color_name, (0.0, 0.0, 0.0))


# ─────────────────────────── 数据加载 ───────────────────────────

def load_page_data(
    page_nums: Optional[list[int]],
    corrections_map: dict,
) -> list[dict]:
    """从 structure/ + text_ocr/ 加载并合并页面数据。"""
    struct_files = sorted(STRUCTURE_DIR.glob("page_*.json"))
    if not struct_files:
        print(f"[错误] {STRUCTURE_DIR} 中没有结构数据，请先运行 structure.py")
        sys.exit(1)

    pages = []
    for sf in struct_files:
        num = int(sf.stem.split("_")[1])
        if page_nums and num not in page_nums:
            continue

        try:
            struct = json.loads(sf.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"  [警告] 跳过 {sf.name}: {e}")
            continue

        ocr_path = OCR_DIR / f"page_{num:03d}.json"
        ocr = {}
        if ocr_path.exists():
            try:
                ocr = json.loads(ocr_path.read_text(encoding="utf-8"))
            except Exception:
                pass

        page = _build_page(num, struct, ocr, corrections_map)
        pages.append(page)

    print(f"[build] 加载 {len(pages)} 页数据")
    return pages


def _build_page(page_num: int, struct: dict, ocr: dict, corrections_map: dict) -> dict:
    """合并 structure + ocr 数据为内部格式。"""
    # 建立 ocr region_id → lines 索引
    ocr_index: dict[str, list] = {}
    for r in ocr.get("text_regions", []):
        ocr_index[r["region_id"]] = r.get("lines", [])

    text_regions = []
    illus_regions = []

    for r in struct.get("regions", []):
        rtype = r.get("type", "")
        rid = r.get("id", "")

        if rtype in ("TEXT_BLOCK", "PAGE_META", "LABEL_SYSTEM", "DIMENSION"):
            lines = ocr_index.get(rid, r.get("lines", []))
            # 应用人工修正
            lines = _apply_corrections(lines, corrections_map)
            confs = [ln.get("confidence", 0.8) for ln in lines if "confidence" in ln]
            avg_conf = sum(confs) / len(confs) if confs else 0.8
            fl = lines[0].get("font_level", "body") if lines else "body"
            color = lines[0].get("color", "black") if lines else "black"
            text_regions.append({
                "id": rid,
                "bbox": r.get("bbox", [0, 0, 100, 20]),
                "font_level": fl,
                "color": color,
                "lines": lines,
                "confidence": avg_conf,
            })

        elif rtype == "ILLUSTRATION":
            illus_regions.append({
                "id": rid,
                "bbox": r.get("bbox", [0, 0, 100, 100]),
                "illustration_type": r.get("illustration_type", "other"),
                "has_overlapping_text": r.get("has_overlapping_text", False),
            })

    all_confs = [r["confidence"] for r in text_regions]
    avg_conf = sum(all_confs) / len(all_confs) if all_confs else None

    orig_img = PAGES_DIR / f"page_{page_num:03d}.png"
    return {
        "page_number": page_num,
        "page_width_px": struct.get("page_width_px", 0),
        "page_height_px": struct.get("page_height_px", 0),
        "text_regions": text_regions,
        "illustration_regions": illus_regions,
        "original_image_path": str(orig_img) if orig_img.exists() else None,
        "average_confidence": avg_conf,
    }


def _apply_corrections(lines: list, corrections_map: dict) -> list:
    if not corrections_map:
        return lines
    result = []
    for ln in lines:
        content = ln.get("content", "")
        for wrong, right in corrections_map.items():
            content = content.replace(wrong, right)
        result.append({**ln, "content": content})
    return result


def _load_corrections(path: Optional[str]) -> dict:
    if not path:
        return {}
    try:
        return json.loads(Path(path).read_text(encoding="utf-8"))
    except Exception as e:
        print(f"  [警告] 修正文件加载失败: {e}")
        return {}

# ─────────────────────────── OCG 支持 ───────────────────────────

def _setup_ocg_support():
    from reportlab.pdfbase import pdfdoc
    if "OCProperties" not in pdfdoc.PDFCatalog.__NoDefault__:
        pdfdoc.PDFCatalog.__NoDefault__.append("OCProperties")
    if "OCProperties" not in pdfdoc.PDFCatalog.__Refs__:
        pdfdoc.PDFCatalog.__Refs__.append("OCProperties")

    _orig_check = pdfdoc.PDFPage.check_format
    def _patched(self, document):
        if hasattr(self, "_ocg_properties"):
            self.Resources.properties = self._ocg_properties
        return _orig_check(self, document)
    pdfdoc.PDFPage.check_format = _patched


class OCGLayer:
    def __init__(self, resource_name, display_name, ref, visible):
        self.resource_name = resource_name
        self.display_name = display_name
        self.ref = ref
        self.visible = visible


class OCGManager:
    def __init__(self, canvas, doc):
        self.canvas = canvas
        self.doc = doc
        self.layers = []
        self._resource_map = {}

    def add_layer(self, resource_name, display_name, visible=True):
        from reportlab.pdfbase import pdfdoc
        ocg_dict = pdfdoc.PDFDictionary({
            "Type": pdfdoc.PDFName("OCG"),
            "Name": pdfdoc.PDFString(display_name),
        })
        ref = self.doc.Reference(ocg_dict)
        layer = OCGLayer(resource_name, display_name, ref, visible)
        self.layers.append(layer)
        self._resource_map[resource_name] = layer
        return layer

    def begin(self, resource_name):
        self.canvas._code.append(f"/OC /{resource_name} BDC")

    def end(self):
        self.canvas._code.append("EMC")

    def attach_to_page(self):
        props = {layer.resource_name: layer.ref for layer in self.layers}
        last_page = self.doc.Pages.pages[-1]
        last_page._ocg_properties = props

    def finalize_catalog(self):
        from reportlab.pdfbase import pdfdoc
        all_refs = [layer.ref for layer in self.layers]
        on_refs  = [layer.ref for layer in self.layers if layer.visible]
        default_cfg = pdfdoc.PDFDictionary({
            "Name":      pdfdoc.PDFString("Default"),
            "BaseState": pdfdoc.PDFName("OFF"),
            "ON":        pdfdoc.PDFArray(on_refs),
            "Order":     pdfdoc.PDFArray([layer.ref for layer in self.layers]),
        })
        self.doc.Catalog.OCProperties = pdfdoc.PDFDictionary({
            "OCGs": pdfdoc.PDFArray(all_refs),
            "D":    default_cfg,
        })


# ─────────────────────────── PDF 工具 ───────────────────────────

_TARGET_PAGE_WIDTH_PT = 595.0


def _compute_page_scale(pw_px, ph_px):
    if pw_px <= 0 or ph_px <= 0:
        return _TARGET_PAGE_WIDTH_PT, 841.89, 1.0, 1.0
    sx = _TARGET_PAGE_WIDTH_PT / pw_px
    return pw_px * sx, ph_px * sx, sx, sx


def draw_colored_text(c, x, y, text, font_name, font_size, default_color):
    from reportlab.pdfbase import pdfmetrics
    pattern = re.compile(r"<(red|blue)>(.*?)</\1>", re.DOTALL)
    parts = []
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            parts.append(("default", text[last:m.start()]))
        parts.append((m.group(1), m.group(2)))
        last = m.end()
    if last < len(text):
        parts.append(("default", text[last:]))
    if not parts:
        parts = [("default", text)]

    cx = x
    for color_key, chunk in parts:
        if not chunk:
            continue
        c.setFillColorRGB(*(get_color_rgb(color_key) if color_key != "default" else default_color))
        c.setFont(font_name, font_size)
        c.drawString(cx, y, chunk)
        cx += pdfmetrics.stringWidth(chunk, font_name, font_size)


def _wrap_text_lines(text, font_name, font_size, max_width):
    from reportlab.pdfbase import pdfmetrics
    raw_lines = text.split("\n")
    result = []
    for raw_line in raw_lines:
        if not raw_line.strip():
            result.append("")
            continue
        clean = re.sub(r"<[^>]+>", "", raw_line)
        if pdfmetrics.stringWidth(clean, font_name, font_size) <= max_width:
            result.append(raw_line)
            continue
        current = ""
        current_clean = ""
        for char in raw_line:
            test_clean = current_clean + char
            if pdfmetrics.stringWidth(test_clean, font_name, font_size) <= max_width:
                current += char
                current_clean = test_clean
            else:
                if current:
                    result.append(current)
                current = char
                current_clean = char
        if current:
            result.append(current)
    return result


def _draw_text_region(c, region, px2pdf, font_map, ph_pt, sx, sy):
    lines_data = region.get("lines", [])
    if not lines_data:
        return
    for line_info in lines_data:
        content = line_info.get("content", "").strip()
        if not content:
            continue
        fl = line_info.get("font_level", "body")
        color_name = line_info.get("color", "black")
        font_name, base_size = font_map.get(fl, font_map["body"])

        # 用 char_height_px 动态计算字号
        char_h_px = line_info.get("char_height_px", 0)
        if char_h_px > 0:
            font_size = char_h_px * sx * 0.75
            font_size = max(6.0, min(font_size, 72.0))
        else:
            font_size = float(base_size)

        y_px = line_info.get("y_px", 0)
        x_px = line_info.get("x_px", 0)
        px, py = px2pdf(x_px, y_px)

        color_rgb = get_color_rgb(color_name)
        draw_colored_text(c, px, py, content, font_name, font_size, color_rgb)


def _draw_svg_layer(c, svg_path, pw_pt, ph_pt):
    """用 svglib 将 SVG 渲染到 canvas。"""
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPDF
        drawing = svg2rlg(str(svg_path))
        if drawing and drawing.width > 0 and drawing.height > 0:
            sx = pw_pt / drawing.width
            sy = ph_pt / drawing.height
            drawing.width  = pw_pt
            drawing.height = ph_pt
            drawing.transform = (sx, 0, 0, sy, 0, 0)
            renderPDF.draw(drawing, c, 0, 0)
    except Exception as e:
        print(f"  [警告] SVG 渲染失败 {svg_path}: {e}")


# ─────────────────────────── PDF 构建 ───────────────────────────

def build_pdf(pages, output_path, font_map):
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.utils import ImageReader

    _setup_ocg_support()
    print(f"[build] 生成 PDF → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    c = rl_canvas.Canvas(str(output_path))
    doc = c._doc
    ocg = OCGManager(c, doc)

    L_SCAN  = ocg.add_layer("L1Scan",  "原始扫描件",  visible=False)
    L_TRACE = ocg.add_layer("L2Trace", "Traced图形",  visible=True)
    L_TEXT  = ocg.add_layer("L3Text",  "数字化文字",  visible=True)
    L_LABEL = ocg.add_layer("L4Label", "标注引线",    visible=True)

    total = len(pages)
    for idx, page in enumerate(pages):
        pnum   = page["page_number"]
        pw_px  = page.get("page_width_px",  908)
        ph_px  = page.get("page_height_px", 1282)
        pw_pt, ph_pt, sx, sy = _compute_page_scale(pw_px, ph_px)
        c.setPageSize((pw_pt, ph_pt))

        def px2pdf(bx, by, bw=None, bh=None):
            px = bx * sx
            py = ph_pt - by * sy
            if bw is not None and bh is not None:
                return px, py - bh * sy, bw * sx, bh * sy
            return px, py

        # Layer 0: 白色背景
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, pw_pt, ph_pt, fill=1, stroke=0)

        # Layer 1: 原始扫描件
        img_path = page.get("original_image_path")
        if img_path and os.path.exists(img_path):
            ocg.begin(L_SCAN.resource_name)
            try:
                c.drawImage(ImageReader(img_path), 0, 0,
                            width=pw_pt, height=ph_pt, preserveAspectRatio=False)
            except Exception as e:
                print(f"  [警告] page_{pnum:03d}: 扫描件读取失败: {e}")
            ocg.end()

        # Layer 2: Traced SVG
        trace_dir = TRACES_DIR / f"page_{pnum:03d}"
        if trace_dir.exists():
            ocg.begin(L_TRACE.resource_name)
            for svg_file in sorted(trace_dir.glob("*.svg")):
                _draw_svg_layer(c, svg_file, pw_pt, ph_pt)
            ocg.end()

        # Layer 3 & 4: 文字
        dim_regions   = [r for r in page["text_regions"] if r.get("font_level") == "dimension"]
        label_regions = [r for r in page["text_regions"] if r.get("font_level") == "label"]
        text_regions  = [r for r in page["text_regions"]
                         if r.get("font_level") not in ("dimension", "label")]

        if text_regions or dim_regions:
            ocg.begin(L_TEXT.resource_name)
            for region in text_regions + dim_regions:
                _draw_text_region(c, region, px2pdf, font_map, ph_pt, sx, sy)
            ocg.end()

        if label_regions:
            ocg.begin(L_LABEL.resource_name)
            for region in label_regions:
                _draw_text_region(c, region, px2pdf, font_map, ph_pt, sx, sy)
            ocg.end()

        c.showPage()
        ocg.attach_to_page()

        if (idx + 1) % 10 == 0 or (idx + 1) == total:
            print(f"  [PDF] {idx+1}/{total} 页完成")

    ocg.finalize_catalog()
    c.save()
    size_mb = output_path.stat().st_size / 1024 / 1024
    print(f"  [PDF] 保存: {output_path} ({size_mb:.1f} MB)")


# ─────────────────────────── PPTX 构建 ───────────────────────────

def build_pptx(pages, output_path, font_map):
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    pptx_font_map = {
        "NotoSansSC-Bold":     "Noto Sans SC",
        "NotoSansSC-Regular":  "Noto Sans SC",
        "NotoSerifSC-Regular": "Noto Serif SC",
        "STHeitiMedium":       "Heiti SC",
        "STHeitiLight":        "Heiti SC",
        "Songti":              "Songti SC",
        "Helvetica":           "Arial",
    }

    print(f"[build] 生成 PPTX → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    total = len(pages)

    for idx, page in enumerate(pages):
        pnum  = page["page_number"]
        pw_px = page.get("page_width_px",  908)
        ph_px = page.get("page_height_px", 1282)

        max_w_emu = int(13.33 * 914400)
        max_h_emu = int(10.0  * 914400)
        aspect = pw_px / ph_px
        if aspect >= 1:
            slide_w = max_w_emu
            slide_h = int(slide_w / aspect)
        else:
            slide_h = max_h_emu
            slide_w = int(slide_h * aspect)

        prs.slide_width  = slide_w
        prs.slide_height = slide_h

        def px2emu(bx, by, bw=None, bh=None):
            ex = int(bx / pw_px * slide_w)
            ey = int(by / ph_px * slide_h)
            if bw is not None and bh is not None:
                return ex, ey, int(bw / pw_px * slide_w), int(bh / ph_px * slide_h)
            return ex, ey

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 插入 traced SVG（转 PNG 后插入）
        trace_dir = TRACES_DIR / f"page_{pnum:03d}"
        if trace_dir.exists():
            for svg_file in sorted(trace_dir.glob("*.svg")):
                try:
                    import cairosvg, io
                    from pptx.util import Emu
                    png_bytes = cairosvg.svg2png(url=str(svg_file),
                                                 output_width=pw_px, output_height=ph_px)
                    slide.shapes.add_picture(
                        io.BytesIO(png_bytes), left=0, top=0,
                        width=slide_w, height=slide_h
                    )
                except Exception as e:
                    print(f"  [警告] page_{pnum:03d}: SVG→PNG 失败 {svg_file.name}: {e}")

        # 文字区域
        for region in page["text_regions"]:
            bbox = region.get("bbox", [0, 0, 200, 40])
            bx, by, bw, bh = bbox
            ex, ey, ew, eh = px2emu(bx, by, bw, bh)
            ew = max(ew, int(0.5 * 914400))
            eh = max(eh, int(0.3 * 914400))

            try:
                txbox = slide.shapes.add_textbox(left=ex, top=ey, width=ew, height=eh)
                tf = txbox.text_frame
                tf.word_wrap = True

                for i, line_info in enumerate(region.get("lines", [])):
                    content = re.sub(r"<[^>]+>", "", line_info.get("content", "")).strip()
                    if not content:
                        continue
                    fl = line_info.get("font_level", "body")
                    color_name = line_info.get("color", "black")
                    rl_font_name, pt_size = font_map.get(fl, font_map["body"])
                    pptx_fname = pptx_font_map.get(rl_font_name, "Microsoft YaHei")

                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    run = para.add_run()
                    run.text = content
                    run.font.name = pptx_fname
                    run.font.size = Pt(pt_size)
                    cr, cg, cb = get_color_rgb(color_name)
                    run.font.color.rgb = RGBColor(int(cr*255), int(cg*255), int(cb*255))
                    if fl in ("title", "subtitle"):
                        run.font.bold = True
            except Exception as e:
                print(f"  [警告] page_{pnum:03d}: 文字框失败: {e}")

        conf = page.get("average_confidence")
        slide.notes_slide.notes_text_frame.text = (
            f"Page {pnum} | Confidence: {conf:.2f}" if conf else f"Page {pnum}"
        )

        if (idx + 1) % 10 == 0 or (idx + 1) == total:
            print(f"  [PPTX] {idx+1}/{total} 页完成")

    prs.save(str(output_path))
    size_mb = output_path.stat().st_size / 1024 / 1024
    print(f"  [PPTX] 保存: {output_path} ({size_mb:.1f} MB)")


# ─────────────────────────── Markdown 构建 ───────────────────────────

def build_markdown(pages, output_path):
    print(f"[build] 生成 Markdown → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    out = [
        "# 教案数字版 — 纯文字",
        "",
        "> 自动生成，请校对。⚠️[?X?]⚠️ 标记为不确定字符。",
        "",
    ]

    for page in pages:
        pnum = page["page_number"]
        out += ["---", "", f"## 第 {pnum} 页", ""]

        conf = page.get("average_confidence")
        if conf:
            out += [f"> 置信度: {int(conf*100)}%", ""]

        # 插图占位
        trace_dir = TRACES_DIR / f"page_{pnum:03d}"
        if trace_dir.exists():
            meta_path = trace_dir / "meta.json"
            if meta_path.exists():
                try:
                    meta = json.loads(meta_path.read_text(encoding="utf-8"))
                    colors = meta.get("colors_found", [])
                    out.append(f"[插图: 颜色通道 {', '.join(colors)}]")
                    out.append("")
                except Exception:
                    pass

        # 文字内容
        for region in page["text_regions"]:
            for line_info in region.get("lines", []):
                content = line_info.get("content", "").strip()
                if not content:
                    continue
                fl = line_info.get("font_level", "body")
                color = line_info.get("color", "black")

                # 格式化
                content = content.replace("[?", "⚠️[?").replace("?]", "?]⚠️")
                if color == "red":
                    content = f"**{content}**（红色）"
                elif color == "blue":
                    content = f"**{content}**（蓝色）"

                if fl == "title":
                    out.append(f"### {content}")
                elif fl == "subtitle":
                    out.append(f"#### {content}")
                else:
                    out.append(content)
            out.append("")

        # 修正记录
        corrections = page.get("corrections", [])
        if corrections:
            out += ["**自动修正：**", ""]
            for c in corrections:
                out.append(f"- `{c.get('original')}` → `{c.get('corrected')}` （{c.get('reason', '')}）")
            out.append("")

    output_path.write_text("\n".join(out), encoding="utf-8")
    print(f"  [MD] 保存: {output_path} ({output_path.stat().st_size/1024:.1f} KB)")


# ─────────────────────────── QA 报告 ───────────────────────────

def build_qa_report(pages, output_path):
    print(f"[build] 生成 QA 报告 → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    total = len(pages)
    confs = [p["average_confidence"] for p in pages if p.get("average_confidence")]
    avg_conf = sum(confs) / len(confs) if confs else None
    all_corrections = []
    for p in pages:
        for c in p.get("corrections", []):
            all_corrections.append({**c, "page": p["page_number"]})

    uncertain_total = sum(
        sum(ln.get("content", "").count("[?") for ln in r.get("lines", []))
        for p in pages for r in p["text_regions"]
    )
    low_conf = [(p["page_number"], p["average_confidence"])
                for p in pages
                if p.get("average_confidence") and p["average_confidence"] < 0.75]

    out = [
        "# 校对报告",
        "",
        "## 汇总统计",
        "",
        f"- 总页数: {total}",
        f"- 平均置信度: {avg_conf:.1%}" if avg_conf else "- 平均置信度: N/A",
        f"- 自动修正次数: {len(all_corrections)}",
        f"- 不确定字符: {uncertain_total}",
        f"- 低置信度页面（<75%）: {len(low_conf)}",
        "",
    ]

    if low_conf:
        out += ["## 低置信度页面", "", "| 页码 | 置信度 |", "|------|--------|"]
        for pnum, conf in sorted(low_conf, key=lambda x: x[1]):
            out.append(f"| 第 {pnum} 页 | {conf:.1%} |")
        out.append("")

    if all_corrections:
        out += ["## 自动修正列表", "", "| 页码 | 原文 | 修正 | 理由 |", "|------|------|------|------|"]
        for c in all_corrections:
            out.append(f"| 第 {c['page']} 页 | `{c.get('original','')}` | `{c.get('corrected','')}` | {c.get('reason','')} |")
        out.append("")

    out += [
        "## 各页详情",
        "",
        "| 页码 | 置信度 | 文字区域数 | 不确定字符 |",
        "|------|--------|-----------|-----------|",
    ]
    for page in pages:
        pnum = page["page_number"]
        conf = page.get("average_confidence")
        conf_str = f"{conf:.1%}" if conf else "N/A"
        text_count = len(page["text_regions"])
        uncertain = sum(
            ln.get("content", "").count("[?")
            for r in page["text_regions"] for ln in r.get("lines", [])
        )
        out.append(f"| 第 {pnum} 页 | {conf_str} | {text_count} | {uncertain} |")

    out += ["", "---", "", "_报告由 build.py 自动生成_", ""]
    output_path.write_text("\n".join(out), encoding="utf-8")
    print(f"  [QA] 保存: {output_path} ({output_path.stat().st_size/1024:.1f} KB)")

    # 同时输出 corrections.json
    corr_path = output_path.parent / "corrections.json"
    if not corr_path.exists():
        corr_map = {c["original"]: c["corrected"] for c in all_corrections}
        corr_path.write_text(json.dumps(corr_map, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"  [QA] corrections.json 已生成，可编辑后用 --corrections 重跑")


# ─────────────────────────── 主流程 ───────────────────────────

def parse_pages_arg(arg: str) -> list[int]:
    pages = set()
    for part in arg.split(","):
        part = part.strip()
        if "-" in part:
            lo, hi = part.split("-", 1)
            pages.update(range(int(lo), int(hi) + 1))
        else:
            pages.add(int(part))
    return sorted(pages)


def main():
    parser = argparse.ArgumentParser(description="Phase 4: 文档合成")
    parser.add_argument("--pages", type=str, default=None)
    parser.add_argument("--format", type=str, default="pdf,pptx,text,qa")
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--download-fonts", action="store_true")
    parser.add_argument("--corrections", type=str, default=None,
                        help="corrections.json 路径")
    args = parser.parse_args()

    formats = {f.strip() for f in args.format.split(",")}
    page_nums = parse_pages_arg(args.pages) if args.pages else None
    corrections_map = _load_corrections(args.corrections)

    pages = load_page_data(page_nums, corrections_map)
    if not pages:
        print("[错误] 没有找到任何页面数据，请先运行前序步骤。")
        sys.exit(1)

    font_map = setup_fonts(download_noto=args.download_fonts)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    pdf_out  = OUTPUT_DIR / "教案_数字版.pdf"
    pptx_out = OUTPUT_DIR / "教案_数字版.pptx"
    md_out   = OUTPUT_DIR / "教案_纯文字.md"
    qa_out   = OUTPUT_DIR / "校对报告.md"

    if "pdf" in formats:
        if not pdf_out.exists() or args.force:
            build_pdf(pages, pdf_out, font_map)
        else:
            print("[build] 跳过 PDF（已存在，使用 --force 覆盖）")

    if "pptx" in formats:
        if not pptx_out.exists() or args.force:
            build_pptx(pages, pptx_out, font_map)
        else:
            print("[build] 跳过 PPTX（已存在，使用 --force 覆盖）")

    if "text" in formats:
        if not md_out.exists() or args.force:
            build_markdown(pages, md_out)
        else:
            print("[build] 跳过 Markdown（已存在，使用 --force 覆盖）")

    if "qa" in formats:
        if not qa_out.exists() or args.force:
            build_qa_report(pages, qa_out)
        else:
            print("[build] 跳过 QA 报告（已存在，使用 --force 覆盖）")

    print()
    print("─" * 50)
    print("[build] 完成！输出目录:", OUTPUT_DIR)
    for p in [pdf_out, pptx_out, md_out, qa_out]:
        if p.exists():
            size = p.stat().st_size
            unit = "MB" if size > 1024*1024 else "KB"
            val  = size/1024/1024 if size > 1024*1024 else size/1024
            print(f"  {p.name}: {val:.1f} {unit}")
    print("─" * 50)


if __name__ == "__main__":
    main()
