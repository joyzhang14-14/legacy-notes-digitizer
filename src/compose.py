"""
Step 4: 合成输出 — PDF(图层切换) + PPT + 纯文字Markdown
用法：
  python src/compose.py                           # 全部页面，全部格式
  python src/compose.py --test 3                  # 前3页
  python src/compose.py --pages 1-5              # 指定范围
  python src/compose.py --format pdf,pptx         # 指定格式
  python src/compose.py --line-style auto         # 线稿版本 strengthen/vectorize/auto
"""

import argparse
import json
import os
import platform
import re
import sys
import time
from datetime import datetime
from pathlib import Path

import yaml
from PIL import Image

ROOT = Path(__file__).parent.parent
CONFIG_PATH = ROOT / "config.yaml"

# ─────────────────────────── 工具函数 ───────────────────────────

def parse_pages_arg(arg: str, total: int) -> list[int]:
    pages = set()
    for part in arg.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            pages.update(range(int(start), int(end) + 1))
        else:
            pages.add(int(part))
    return sorted(p for p in pages if 1 <= p <= total)


def px_to_pdf(bbox_px: list, img_w: int, img_h: int,
              page_w: float, page_h: float) -> tuple:
    """像素坐标（左上原点，Y向下）→ PDF坐标（左下原点，Y向上）"""
    x, y, w, h = bbox_px
    sx = page_w / img_w
    sy = page_h / img_h
    return (x * sx,
            page_h - (y + h) * sy,
            w * sx,
            h * sy)


def get_image_size(path: Path) -> tuple[int, int]:
    with Image.open(path) as img:
        return img.size  # (width, height)


# ─────────────────────────── 字体注册 ───────────────────────────

def register_fonts() -> dict[str, str]:
    """注册中文字体，返回 {逻辑名: 已注册名} 映射"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    sys_name = platform.system()
    candidates = {
        "Darwin": [
            ("STHeiti",   "/System/Library/Fonts/STHeiti Medium.ttc"),
            ("Songti",    "/System/Library/Fonts/Supplemental/Songti.ttc"),
            ("STHeitiLight", "/System/Library/Fonts/STHeiti Light.ttc"),
        ],
        "Linux": [
            ("NotoSansCJK", "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
            ("WQYMicroHei", "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc"),
        ],
        "Windows": [
            ("SimHei",   "C:/Windows/Fonts/simhei.ttf"),
            ("SimSun",   "C:/Windows/Fonts/simsun.ttc"),
            ("KaiTi",    "C:/Windows/Fonts/simkai.ttf"),
            ("FangSong", "C:/Windows/Fonts/simfang.ttf"),
        ],
    }

    registered: dict[str, str] = {}
    for name, path in candidates.get(sys_name, candidates["Darwin"]):
        if not Path(path).exists():
            continue
        try:
            pdfmetrics.registerFont(TTFont(name, path))
            registered[name] = path
        except Exception:
            pass

    if not registered:
        print("[警告] 未找到中文字体，文字可能显示为方块")
    return registered


def pick_font(font_level: str, registered: dict[str, str]) -> str:
    """选择字体名（带 fallback）"""
    # 有哪些字体就用哪些；macOS 上 STHeiti 覆盖所有等级
    preference = {
        "title":      ["STHeiti", "SimHei", "NotoSansCJK"],
        "subtitle":   ["STHeiti", "SimHei", "NotoSansCJK"],
        "body":       ["Songti",  "STHeiti", "SimSun", "NotoSansCJK"],
        "annotation": ["Songti",  "STHeiti", "KaiTi",  "NotoSansCJK"],
        "label":      ["STHeiti", "SimHei",  "NotoSansCJK"],
    }
    for name in preference.get(font_level, ["STHeiti"]):
        if name in registered:
            return name
    # 最终 fallback
    return next(iter(registered), "Helvetica")


def calc_font_size(region: dict, size_range: list,
                   bbox_w_pt: float, bbox_h_pt: float) -> float:
    """根据 bbox 高度估算合适字号，限制在 size_range 内"""
    content = region.get("content", "")
    lines = max(1, content.count("\n") + 1)
    est = bbox_h_pt / (lines * 1.4)
    lo, hi = size_range
    return max(lo, min(hi, est))


# ─────────────────────────── 数据加载 ───────────────────────────

def load_page_data(page_num: int, data_dir: Path,
                   pages_dir: Path, enhanced_dir: Path,
                   line_style: str, cfg: dict) -> dict | None:
    """加载单页所有数据，组装成统一结构，返回 None 表示跳过"""
    json_path = data_dir / f"page_{page_num:03d}.json"
    if not json_path.exists():
        return None
    try:
        page = json.loads(json_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError) as e:
        print(f"  [错误] page_{page_num:03d}.json: {e}")
        return None

    orig_path = pages_dir / f"page_{page_num:03d}.png"
    if not orig_path.exists():
        return None

    img_w = page.get("page_width_px") or 0
    img_h = page.get("page_height_px") or 0
    if not img_w or not img_h:
        img_w, img_h = get_image_size(orig_path)
        page["page_width_px"] = img_w
        page["page_height_px"] = img_h

    # 决定每个插图用哪个版本的线稿
    meta_path = enhanced_dir / f"page_{page_num:03d}_meta.json"
    meta = {}
    if meta_path.exists():
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
        except Exception:
            pass

    id_to_rec = {ill["id"]: ill.get("recommended_version", "strengthen")
                 for ill in meta.get("illustrations", [])}

    for ill in page.get("illustration_regions", []):
        iid = ill["id"]
        if line_style == "auto":
            version = id_to_rec.get(iid, "strengthen")
        else:
            version = line_style

        enhanced_png = enhanced_dir / f"page_{page_num:03d}_{version}.png"
        if not enhanced_png.exists():
            # fallback：另一个版本
            other = "vectorize" if version == "strengthen" else "strengthen"
            fallback = enhanced_dir / f"page_{page_num:03d}_{other}.png"
            enhanced_png = fallback if fallback.exists() else None

        ill["enhanced_path"] = str(enhanced_png) if enhanced_png else None

    page["original_path"] = str(orig_path)
    page["page_number"] = page_num
    return page


# ─────────────────────────── PDF 生成 ───────────────────────────

def _patch_catalog_for_ocg():
    """让 reportlab PDFCatalog 支持 OCProperties 字段"""
    from reportlab.pdfbase import pdfdoc
    if "OCProperties" not in pdfdoc.PDFCatalog.__NoDefault__:
        pdfdoc.PDFCatalog.__NoDefault__.append("OCProperties")
    if "OCProperties" not in pdfdoc.PDFCatalog.__Refs__:
        pdfdoc.PDFCatalog.__Refs__.append("OCProperties")


def build_pdf(pages_data: list[dict], output_path: Path,
              registered_fonts: dict, cfg: dict) -> None:
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.pdfbase import pdfdoc
    from reportlab.pdfbase.pdfdoc import (PDFDictionary, PDFArray,
                                           PDFName, PDFString,
                                           PDFResourceDictionary)

    _patch_catalog_for_ocg()

    DPI = cfg["preprocessing"]["dpi"]
    mask_padding = cfg["compose"]["pdf"]["mask_padding"]
    colors_cfg = cfg["compose"]["colors"]
    fonts_cfg = cfg["compose"]["fonts"]

    c = rl_canvas.Canvas(str(output_path), pageCompression=1)
    doc = c._doc

    # ── 创建 OCG（可切换图层）──
    ocg_text = PDFDictionary({
        "Type": PDFName("OCG"),
        "Name": PDFString("数字化文字"),
        "Usage": PDFDictionary({"Print": PDFDictionary({"PrintState": PDFName("ON")})}),
    })
    ocg_lines = PDFDictionary({
        "Type": PDFName("OCG"),
        "Name": PDFString("增强线稿"),
        "Usage": PDFDictionary({"Print": PDFDictionary({"PrintState": PDFName("ON")})}),
    })
    ocg_text_ref = doc.Reference(ocg_text, "OCGText")
    ocg_lines_ref = doc.Reference(ocg_lines, "OCGLines")

    on_arr = PDFArray([ocg_text_ref, ocg_lines_ref])
    default_d = PDFDictionary({
        "BaseState": PDFName("ON"),
        "ON": on_arr,
        "OFF": PDFArray([]),
        "Order": on_arr,
    })
    doc.Catalog.OCProperties = PDFDictionary({"OCGs": on_arr, "D": default_d})

    for page in pages_data:
        img_w = page["page_width_px"]
        img_h = page["page_height_px"]
        page_w = img_w / DPI * 72.0
        page_h = img_h / DPI * 72.0
        c.setPageSize((page_w, page_h))

        # ── Layer 1: 原始扫描件底图（始终可见）──
        c.drawImage(page["original_path"], 0, 0, width=page_w, height=page_h,
                    preserveAspectRatio=False)

        # ── Layer 3: 增强线稿（可切换）──
        c._code.append("/OC /LyrLines BDC")
        for ill in page.get("illustration_regions", []):
            ep = ill.get("enhanced_path")
            if ep and Path(ep).exists():
                bx, by, bw, bh = px_to_pdf(ill["bbox"], img_w, img_h, page_w, page_h)
                c.drawImage(ep, bx, by, width=bw, height=bh,
                            mask="auto", preserveAspectRatio=False)
        c._code.append("EMC")

        # ── Layer 4: 数字化文字 + 白色遮罩（可切换）──
        c._code.append("/OC /LyrText BDC")
        for region in page.get("text_regions", []):
            bx, by, bw, bh = px_to_pdf(region["bbox"], img_w, img_h, page_w, page_h)
            pad = mask_padding

            # 白色遮罩（覆盖原手写文字）
            c.setFillColorRGB(1, 1, 1)
            c.rect(bx - pad, by - pad, bw + 2 * pad, bh + 2 * pad,
                   fill=1, stroke=0)

            # 绘制数字化文字
            color_key = region.get("color", "black")
            rgb = [v / 255.0 for v in colors_cfg.get(color_key, [0, 0, 0])]
            c.setFillColorRGB(*rgb)

            font_level = region.get("font_level", "body")
            font_name = pick_font(font_level, registered_fonts)
            font_size = calc_font_size(region,
                                       fonts_cfg[font_level]["size_range"],
                                       bw, bh)
            c.setFont(font_name, font_size)
            leading = font_size * 1.4

            content = region.get("content", "")
            text_obj = c.beginText(bx, by + bh - font_size)
            text_obj.setFont(font_name, font_size)
            text_obj.setLeading(leading)
            for line in content.split("\n"):
                text_obj.textLine(line)
            c.drawText(text_obj)
        c._code.append("EMC")

        c.showPage()

        # 给刚生成的页面注入 OCG Properties（需要在 showPage 之后修改）
        last_page = doc.Pages.pages[-1]
        if not last_page.Resources:
            last_page.Resources = PDFResourceDictionary()
        last_page.Resources.Properties = {
            "LyrText": ocg_text_ref,
            "LyrLines": ocg_lines_ref,
        }

    c.save()


# ─────────────────────────── PPTX 生成 ───────────────────────────

def build_pptx(pages_data: list[dict], output_path: Path,
               registered_fonts: dict, cfg: dict) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    pptx_cfg = cfg["compose"]["pptx"]
    fonts_cfg = cfg["compose"]["fonts"]
    colors_cfg = cfg["compose"]["colors"]
    DPI = cfg["preprocessing"]["dpi"]
    qa_threshold = cfg.get("qa", {}).get("low_confidence_threshold", 0.7)

    prs = Presentation()
    prs.slide_width = Inches(pptx_cfg["slide_width"])
    prs.slide_height = Inches(pptx_cfg["slide_height"])

    blank_layout = prs.slide_layouts[6]  # blank layout

    for page in pages_data:
        img_w = page["page_width_px"]
        img_h = page["page_height_px"]
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        def px_to_emu(bbox):
            x, y, w, h = bbox
            ex = int(x / img_w * slide_w)
            ey = int(y / img_h * slide_h)
            ew = int(w / img_w * slide_w)
            eh = int(h / img_h * slide_h)
            return ex, ey, ew, eh

        slide = prs.slides.add_slide(blank_layout)

        # ── 背景：原始扫描件 ──
        slide.shapes.add_picture(
            page["original_path"],
            Emu(0), Emu(0), width=slide_w, height=slide_h,
        )

        # ── 增强线稿（插图区域替换）──
        for ill in page.get("illustration_regions", []):
            ep = ill.get("enhanced_path")
            if ep and Path(ep).exists():
                ex, ey, ew, eh = px_to_emu(ill["bbox"])
                slide.shapes.add_picture(ep, Emu(ex), Emu(ey),
                                         width=Emu(ew), height=Emu(eh))

        # ── 文字区域：白色遮罩 + 文本框 ──
        for region in page.get("text_regions", []):
            ex, ey, ew, eh = px_to_emu(region["bbox"])

            # 白色矩形遮罩
            from pptx.util import Pt as _Pt
            mask = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
                Emu(ex - 4), Emu(ey - 4), Emu(ew + 8), Emu(eh + 8),
            )
            mask.fill.solid()
            mask.fill.fore_color.rgb = RGBColor(255, 255, 255)
            mask.line.fill.background()

            # 文本框
            txbox = slide.shapes.add_textbox(
                Emu(ex), Emu(ey), Emu(ew), Emu(eh),
            )
            tf = txbox.text_frame
            tf.word_wrap = True

            color_key = region.get("color", "black")
            rgb_vals = colors_cfg.get(color_key, [0, 0, 0])
            font_level = region.get("font_level", "body")
            font_name = fonts_cfg[font_level].get("name", "STHeiti")
            font_size_range = fonts_cfg[font_level]["size_range"]
            font_size = int(calc_font_size(region, font_size_range,
                                           ew / 9144, eh / 9144))  # emu→inch 粗略

            content = region.get("content", "")
            for i, line in enumerate(content.split("\n")):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = para.add_run()
                run.text = line
                run.font.size = Pt(max(8, font_size))
                run.font.color.rgb = RGBColor(*rgb_vals)
                run.font.name = font_name
                if font_level == "title":
                    run.font.bold = True

        # ── 备注：页码 + 置信度 ──
        confs = [r.get("confidence", 1.0) for r in page.get("text_regions", [])]
        avg_conf = sum(confs) / len(confs) if confs else 1.0
        low_conf = avg_conf < qa_threshold
        note_text = (
            f"第 {page['page_number']} 页 | "
            f"类型: {page.get('page_type', '?')} | "
            f"平均置信度: {avg_conf:.2f}"
        )
        if low_conf:
            note_text += " ⚠️ 建议人工校对"
        slide.notes_slide.notes_text_frame.text = note_text

    prs.save(str(output_path))


# ─────────────────────────── Markdown 生成 ───────────────────────────

def build_markdown(pages_data: list[dict], output_path: Path,
                   cfg: dict) -> None:
    qa_threshold = cfg.get("qa", {}).get("low_confidence_threshold", 0.7)
    total = len(pages_data)
    today = datetime.now().strftime("%Y-%m-%d")

    lines = [
        "# 教案文字记录",
        "",
        f"> 生成时间: {today}",
        f"> 总页数: {total}",
        "> 工具: Legacy Notes Digitizer",
        "",
        "---",
        "",
    ]

    for page in pages_data:
        num = page["page_number"]
        page_type = page.get("page_type", "")
        lines.append(f"## 第 {num} 页")
        lines.append("")

        if page_type == "illustration_only":
            lines.append("_（本页为纯插图）_")
            lines.append("")

        text_regions = page.get("text_regions", [])
        ill_regions = page.get("illustration_regions", [])

        # 按 Y 坐标排序（从上到下）
        text_regions = sorted(text_regions, key=lambda r: r["bbox"][1])

        # 先输出文字区域
        ill_counter = 0
        ill_used = set()

        for region in text_regions:
            content = region.get("content", "").strip()
            if not content:
                continue

            font_level = region.get("font_level", "body")
            color = region.get("color", "black")
            conf = region.get("confidence", 1.0)

            # 低置信度标记
            if conf < qa_threshold:
                content = f"⚠️ {content} ⚠️"

            # 颜色标注
            color_suffix = ""
            if color == "red":
                color_suffix = "（红色标注）"
            elif color == "blue":
                color_suffix = "（蓝色标注）"

            # 按层级格式化
            if font_level == "title":
                lines.append(f"**{content}**{color_suffix}")
            elif font_level == "subtitle":
                lines.append(f"### {content}{color_suffix}")
            elif font_level == "annotation":
                lines.append(f"> {content}{color_suffix}")
            elif font_level == "label":
                lines.append(f"**标注文字：** {content}{color_suffix}")
            else:  # body
                lines.append(f"{content}{color_suffix}")

            lines.append("")

        # 插图区域
        for ill in ill_regions:
            ill_counter += 1
            desc = ill.get("description", "")
            ill_type = ill.get("type", "other")
            lines.append(f"[插图 {ill_counter}: {ill_type} — {desc}]")
            lines.append("")

        lines.append("---")
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")


# ─────────────────────────── 主流程 ───────────────────────────

def main():
    parser = argparse.ArgumentParser(description="合成输出：PDF / PPT / Markdown")
    parser.add_argument("--test", type=int, metavar="N", help="只处理前 N 页")
    parser.add_argument("--pages", type=str, metavar="RANGE",
                        help="页码范围，如 '1-5' 或 '1,3,7'")
    parser.add_argument("--format", type=str, default="all",
                        help="输出格式，逗号分隔：pdf,pptx,text（默认 all）")
    parser.add_argument("--line-style",
                        choices=["strengthen", "vectorize", "auto"],
                        default="auto", help="线稿版本（默认 auto）")
    parser.add_argument("--corrections", type=str,
                        help="修正后的 JSON 目录（可选）")
    args = parser.parse_args()

    cfg = yaml.safe_load(CONFIG_PATH.read_text(encoding="utf-8"))
    inter_dir = ROOT / cfg["paths"]["intermediate_dir"]
    out_dir = ROOT / cfg["paths"]["output_dir"]
    pages_dir = inter_dir / "pages"
    data_dir = inter_dir / "page_data"
    enhanced_dir = inter_dir / "enhanced_lines"
    out_dir.mkdir(parents=True, exist_ok=True)

    # 确定要处理的页码
    all_jsons = sorted(data_dir.glob("page_*.json"))
    total = len(all_jsons)
    if total == 0:
        print(f"[错误] {data_dir} 中没有 JSON，请先运行 analyze.py")
        sys.exit(1)

    if args.pages:
        target_nums = parse_pages_arg(args.pages, total)
    elif args.test:
        target_nums = list(range(1, args.test + 1))
        print(f"[测试模式] 只处理前 {args.test} 页")
    else:
        target_nums = list(range(1, total + 1))

    # 确定输出格式
    fmt_input = args.format.lower()
    formats = {"pdf", "pptx", "text"} if fmt_input == "all" \
        else {f.strip() for f in fmt_input.split(",")}

    print(f"处理 {len(target_nums)} 页，格式: {', '.join(sorted(formats))}，"
          f"线稿: {args.line_style}")

    # 加载所有页面数据
    start = time.time()
    pages_data = []
    errors = []
    for i, num in enumerate(target_nums):
        page = load_page_data(num, data_dir, pages_dir, enhanced_dir,
                              args.line_style, cfg)
        if page is None:
            errors.append(num)
            continue
        pages_data.append(page)
        if (i + 1) % 10 == 0 or (i + 1) == len(target_nums):
            print(f"  [{i + 1}/{len(target_nums)}] 数据加载中...")

    print(f"成功加载 {len(pages_data)} 页，{len(errors)} 页缺失数据")
    if not pages_data:
        print("[错误] 没有可用的页面数据")
        sys.exit(1)

    # 注册字体（PDF/PPT 共用）
    registered_fonts = {}
    if "pdf" in formats or "pptx" in formats:
        registered_fonts = register_fonts()
        if registered_fonts:
            print(f"已注册字体：{list(registered_fonts.keys())}")

    # ── 生成 PDF ──
    if "pdf" in formats:
        pdf_path = out_dir / "教案_数字版.pdf"
        print(f"\n生成 PDF: {pdf_path.name} ...")
        t0 = time.time()
        try:
            build_pdf(pages_data, pdf_path, registered_fonts, cfg)
            size_mb = pdf_path.stat().st_size / 1024 / 1024
            print(f"  完成 ({time.time()-t0:.1f}s, {size_mb:.1f} MB)")
        except Exception as e:
            import traceback
            print(f"  [错误] PDF 生成失败: {e}")
            traceback.print_exc()

    # ── 生成 PPTX ──
    if "pptx" in formats:
        pptx_path = out_dir / "教案_数字版.pptx"
        print(f"\n生成 PPTX: {pptx_path.name} ...")
        t0 = time.time()
        try:
            build_pptx(pages_data, pptx_path, registered_fonts, cfg)
            size_mb = pptx_path.stat().st_size / 1024 / 1024
            print(f"  完成 ({time.time()-t0:.1f}s, {size_mb:.1f} MB)")
        except Exception as e:
            import traceback
            print(f"  [错误] PPTX 生成失败: {e}")
            traceback.print_exc()

    # ── 生成 Markdown ──
    if "text" in formats:
        md_path = out_dir / "教案_纯文字.md"
        print(f"\n生成 Markdown: {md_path.name} ...")
        t0 = time.time()
        try:
            build_markdown(pages_data, md_path, cfg)
            size_kb = md_path.stat().st_size / 1024
            print(f"  完成 ({time.time()-t0:.1f}s, {size_kb:.1f} KB)")
        except Exception as e:
            import traceback
            print(f"  [错误] Markdown 生成失败: {e}")
            traceback.print_exc()

    # ── 汇总 ──
    print(f"\n{'='*50}")
    print(f"合成完成，总耗时 {time.time()-start:.1f}s")
    for name in ["教案_数字版.pdf", "教案_数字版.pptx", "教案_纯文字.md"]:
        p = out_dir / name
        if p.exists():
            size = p.stat().st_size
            unit = "KB" if size < 1024 * 1024 else "MB"
            val = size / 1024 if size < 1024 * 1024 else size / 1024 / 1024
            print(f"  {name:<25} {val:.1f} {unit}")
    print(f"  输出目录: {out_dir}")


if __name__ == "__main__":
    main()
