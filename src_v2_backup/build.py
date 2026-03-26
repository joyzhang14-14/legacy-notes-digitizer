"""
Phase 3: 从白纸创建全新文档 — PDF(5层) + PPT + 纯文字

用法：
  python src/build.py                                    # 生成全部输出
  python src/build.py --pages 1-3                       # 只处理第1-3页
  python src/build.py --format pdf,pptx,text            # 指定输出格式
  python src/build.py --line-style strengthen           # 指定插图样式
  python src/build.py --force                           # 强制覆盖已有输出

数据来源：
  intermediate/page_data/*.json  — 结构+OCR数据（支持 page_data 和 text_ocr 两种格式）
  intermediate/pages/*.png       — 原始扫描件图片
  intermediate/illustrations/    — 增强插图（如有）
"""

import argparse
import json
import math
import os
import re
import sys
import urllib.request
from pathlib import Path
from typing import Optional

# ─────────────────────────── 路径 ───────────────────────────

ROOT = Path(__file__).parent.parent
PAGES_DIR       = ROOT / "intermediate" / "pages"
PAGE_DATA_DIR   = ROOT / "intermediate" / "page_data"
ILLUSTRATIONS_DIR = ROOT / "intermediate" / "illustrations"
OUTPUT_DIR      = ROOT / "output"
FONTS_DIR       = ROOT / "fonts"

# ─────────────────────────── 字体 ───────────────────────────

# Noto CJK 字体下载 URL（用于跨平台嵌入）
_NOTO_URLS = {
    "NotoSansSC-Regular": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTF/SimplifiedChinese/"
        "NotoSansCJKsc-Regular.otf"
    ),
    "NotoSansSC-Bold": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTF/SimplifiedChinese/"
        "NotoSansCJKsc-Bold.otf"
    ),
    "NotoSerifSC-Regular": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Serif/OTF/SimplifiedChinese/"
        "NotoSerifCJKsc-Regular.otf"
    ),
    "NotoSerifSC-Bold": (
        "https://github.com/googlefonts/noto-cjk/raw/main/Serif/OTF/SimplifiedChinese/"
        "NotoSerifCJKsc-Bold.otf"
    ),
}

# macOS 系统字体备选（按优先级排列）
_MACOS_FALLBACKS = [
    # (reportlab名, 文件路径, subfontIndex, 类型: "sans"|"serif")
    ("STHeitiMedium", "/System/Library/Fonts/STHeiti Medium.ttc", 0, "sans-bold"),
    ("STHeitiLight",  "/System/Library/Fonts/STHeiti Light.ttc",  0, "sans"),
    ("Songti",        "/System/Library/Fonts/Supplemental/Songti.ttc", 1, "serif"),
]

# 当前会话中注册成功的字体（名称 → 是否成功）
_REGISTERED_FONTS: dict[str, str] = {}  # reportlab名 → 文件路径


def _try_register(rl_name: str, path: str, subfont_index: int = 0) -> bool:
    """尝试注册一个字体文件，失败静默返回 False。"""
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    try:
        pdfmetrics.registerFont(TTFont(rl_name, path, subfontIndex=subfont_index))
        _REGISTERED_FONTS[rl_name] = path
        return True
    except Exception as e:
        print(f"  [字体] 注册 {rl_name} 失败: {e}")
        return False


def _download_noto() -> dict[str, str]:
    """下载 Noto CJK 字体到 fonts/ 目录。返回 {name: path} 字典。"""
    FONTS_DIR.mkdir(parents=True, exist_ok=True)
    downloaded = {}
    for name, url in _NOTO_URLS.items():
        path = FONTS_DIR / f"{name}.otf"
        if not path.exists():
            print(f"  [字体] 下载 {name}...")
            try:
                urllib.request.urlretrieve(url, str(path))
                print(f"  [字体] ✓ {name} 下载完成")
                downloaded[name] = str(path)
            except Exception as e:
                print(f"  [字体] ✗ {name} 下载失败: {e}")
        else:
            downloaded[name] = str(path)
    return downloaded


def setup_fonts(download_noto: bool = False) -> dict:
    """
    注册 CJK 字体。优先使用 fonts/ 目录中的 Noto CJK，
    如无则用 macOS 系统字体。
    返回字体映射 dict: {逻辑名 → reportlab字体名}
    """
    print("[build] 初始化字体...")

    # 1. 尝试 fonts/ 目录中的 Noto CJK
    noto_available = {}
    for name in _NOTO_URLS:
        path = FONTS_DIR / f"{name}.otf"
        if path.exists():
            if _try_register(name, str(path)):
                noto_available[name] = str(path)

    # 2. 如果需要下载且没有 Noto 字体
    if download_noto and len(noto_available) < 4:
        dl = _download_noto()
        for name, path in dl.items():
            if name not in noto_available:
                if _try_register(name, path):
                    noto_available[name] = path

    # 3. 安装下载的 Noto 到系统（用于 PPTX）
    if noto_available:
        _install_noto_for_pptx(noto_available)

    # 4. macOS 系统字体兜底
    system_fonts = {}
    for rl_name, path, idx, ftype in _MACOS_FALLBACKS:
        if os.path.exists(path):
            if _try_register(rl_name, path, idx):
                system_fonts[ftype] = rl_name

    # 5. 构建逻辑映射
    def pick(noto_key: str, fallback_type: str) -> str:
        if noto_key in noto_available:
            return noto_key
        return system_fonts.get(fallback_type,
               system_fonts.get("sans", "Helvetica"))

    font_map = {
        "title":       (pick("NotoSansSC-Bold",     "sans-bold"), 28),
        "subtitle":    (pick("NotoSansSC-Bold",     "sans-bold"), 20),
        "body":        (pick("NotoSerifSC-Regular", "serif"),     14),
        "annotation":  (pick("NotoSerifSC-Regular", "serif"),     12),
        "label":       (pick("NotoSansSC-Regular",  "sans"),      11),
        "dimension":   (pick("NotoSansSC-Regular",  "sans"),      10),
    }

    print(f"  [字体] 映射: title={font_map['title'][0]}, body={font_map['body'][0]}")
    return font_map


def _install_noto_for_pptx(noto_paths: dict[str, str]) -> None:
    """将 Noto 字体安装到 ~/Library/Fonts/（macOS用户级，PPTX用）。"""
    import platform, shutil
    if platform.system() != "Darwin":
        return
    user_fonts = Path.home() / "Library" / "Fonts"
    user_fonts.mkdir(parents=True, exist_ok=True)
    for name, src in noto_paths.items():
        dst = user_fonts / Path(src).name
        if not dst.exists():
            try:
                shutil.copy2(src, dst)
                print(f"  [字体] 系统安装: {name}")
            except Exception:
                pass


# ─────────────────────────── 颜色 ───────────────────────────

_COLOR_RGB = {
    "black": (0.0,  0.0,  0.0),
    "red":   (0.80, 0.08, 0.08),
    "blue":  (0.10, 0.20, 0.72),
}

def get_color_rgb(color_name: str) -> tuple[float, float, float]:
    return _COLOR_RGB.get(color_name, (0.0, 0.0, 0.0))


# ─────────────────────────── 数据加载 ─────────────────���─────────

def load_page_data(page_nums: Optional[list[int]] = None) -> list[dict]:
    """
    加载 page_data/*.json（page_data 格式）。
    支持两种格式：
      - page_data 格式: text_regions[].content, .font_level, .color
      - text_ocr 格式:  text_regions[].lines[].content, .font_level
    统一转换成内部格式后返回。
    """
    data_dir = PAGE_DATA_DIR
    if not data_dir.exists():
        print(f"[错误] 找不到数据目录: {data_dir}")
        sys.exit(1)

    json_files = sorted(data_dir.glob("page_*.json"))
    if not json_files:
        print(f"[错误] {data_dir} 中没有页面数据")
        sys.exit(1)

    pages = []
    for jf in json_files:
        if jf.name.endswith(".error.txt"):
            continue
        num = int(jf.stem.split("_")[1])
        if page_nums and num not in page_nums:
            continue

        try:
            raw = json.loads(jf.read_text(encoding="utf-8"))
        except Exception as e:
            print(f"  [警告] 跳过 {jf.name}: {e}")
            continue

        page = _normalize_page(raw, num)
        pages.append(page)

    print(f"[build] 加载 {len(pages)} 页数据")
    return pages


def _normalize_page(raw: dict, page_num: int) -> dict:
    """
    将原始 JSON 统一转换为内部格式。
    内部格式 text_region:
      {id, bbox:[x,y,w,h], lines:[{content, font_level, color, confidence}], font_level, color}
    """
    text_regions = []
    for r in raw.get("text_regions", []):
        # 判断格式
        if "lines" in r:
            # text_ocr 格式
            lines = r["lines"]
            merged_content = "\n".join(ln.get("content", "") for ln in lines)
            font_level = lines[0].get("font_level", "body") if lines else "body"
            color = lines[0].get("color", "black") if lines else "black"
        else:
            # page_data 格式
            content = r.get("content", "")
            font_level = r.get("font_level", "body")
            color = r.get("color", "black")
            confidence = r.get("confidence", 0.8)
            lines = [{"content": content, "font_level": font_level,
                      "color": color, "confidence": confidence}]
            merged_content = content

        text_regions.append({
            "id": r.get("id", ""),
            "bbox": r.get("bbox", [0, 0, 100, 20]),
            "font_level": font_level,
            "color": color,
            "content": merged_content,
            "lines": lines,
            "confidence": r.get("confidence",
                                lines[0].get("confidence", 0.8) if lines else 0.8),
        })

    # 插图区域
    illus_regions = []
    for r in raw.get("illustration_regions", []):
        illus_regions.append({
            "id": r.get("id", ""),
            "bbox": r.get("bbox", [0, 0, 100, 100]),
            "illustration_type": r.get("illustration_type", "other"),
            "has_overlapping_text": r.get("has_overlapping_text", False),
        })

    # 图片路径
    orig_img = PAGES_DIR / f"page_{page_num:03d}.png"

    return {
        "page_number": page_num,
        "page_type": raw.get("page_type", "text_only"),
        "page_width_px": raw.get("page_width_px", 0),
        "page_height_px": raw.get("page_height_px", 0),
        "text_regions": text_regions,
        "illustration_regions": illus_regions,
        "original_image_path": str(orig_img) if orig_img.exists() else None,
        "average_confidence": raw.get("average_confidence",
                                      _mean([r["confidence"] for r in text_regions])),
        "notes": raw.get("notes", ""),
    }


def _mean(vals: list) -> Optional[float]:
    vals = [v for v in vals if v is not None]
    return sum(vals) / len(vals) if vals else None


# ─────────────────────────── 插图路径解析 ───────────────────────────

def get_illustration_path(page_num: int, region_id: str, style: str) -> Optional[str]:
    """
    根据线稿样式，返回插图文件路径。
    style: "original" | "strengthen" | "vectorize" | "inpainted" | "auto"
    """
    base_dir = ILLUSTRATIONS_DIR / f"page_{page_num:03d}"
    if not base_dir.exists():
        return None

    meta_path = base_dir / f"{region_id}_meta.json"
    if not meta_path.exists():
        return None

    try:
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
    except Exception:
        return None

    versions = meta.get("versions", {})

    if style == "auto":
        style = meta.get("recommended_version", "strengthen")

    # 检查是否有 overlapping_text，优先用 inpainted
    if meta.get("has_overlapping_text") and "inpainted" in versions:
        if style not in ("original",):
            style = "inpainted"

    # 样式到文件名的映射
    style_key_map = {
        "strengthen": "strengthened",
        "vectorize":  "vectorized_png",
        "inpainted":  "inpainted",
        "original":   "original",
    }
    filename = versions.get(style_key_map.get(style, style)) or versions.get("strengthened") or versions.get("original")
    if not filename:
        return None

    path = base_dir / filename
    return str(path) if path.exists() else None


# ─────────────────────────── OCG 支持 ───────────────────────────

def _setup_ocg_support():
    """
    在 reportlab 的 PDFCatalog 和 PDFPage 中注入 OCG 支持。
    只需调用一次。
    """
    from reportlab.pdfbase import pdfdoc

    # 让 PDFCatalog.format() 能输出 OCProperties（作为间接引用）
    if "OCProperties" not in pdfdoc.PDFCatalog.__NoDefault__:
        pdfdoc.PDFCatalog.__NoDefault__.append("OCProperties")
    if "OCProperties" not in pdfdoc.PDFCatalog.__Refs__:
        pdfdoc.PDFCatalog.__Refs__.append("OCProperties")

    # 让 PDFPage.check_format() 能输出 Properties（OCG资源）
    _orig = pdfdoc.PDFPage.check_format

    def _patched(self, document):
        _orig(self, document)
        ocg_props = getattr(self, "_ocg_properties", None)
        if ocg_props and isinstance(self.Resources, pdfdoc.PDFResourceDictionary):
            self.Resources.Properties = ocg_props

    pdfdoc.PDFPage.check_format = _patched


class OCGLayer:
    """代表一个 OCG 图层。"""
    def __init__(self, resource_name: str, display_name: str, ref, visible: bool):
        self.resource_name = resource_name
        self.display_name = display_name
        self.ref = ref
        self.visible = visible


class OCGManager:
    """管理 PDF 中的 OCG 图层。"""

    def __init__(self, canvas, doc):
        self.canvas = canvas
        self.doc = doc
        self.layers: list[OCGLayer] = []
        self._resource_map: dict[str, OCGLayer] = {}

    def add_layer(self, resource_name: str, display_name: str, visible: bool = True) -> OCGLayer:
        from reportlab.pdfbase import pdfdoc
        ocg_dict = pdfdoc.PDFDictionary({
            "Type": pdfdoc.PDFName("OCG"),
            "Name": pdfdoc.PDFString(display_name),
        })
        ref = self.doc.Reference(ocg_dict)
        layer = OCGLayer(resource_name, display_name, ref, visible)
        self.layers.append(layer)
        self._resource_map[resource_name] = layer
        return layer

    def begin(self, resource_name: str):
        """在画布流中插入 BDC 标记。"""
        self.canvas._code.append(f"/OC /{resource_name} BDC")

    def end(self):
        """在画布流中插入 EMC 标记。"""
        self.canvas._code.append("EMC")

    def attach_to_page(self):
        """
        将 OCG 资源映射注入到刚刚 showPage() 添加的最后一页。
        必须在 showPage() 之后立即调用。
        """
        props = {layer.resource_name: layer.ref for layer in self.layers}
        last_page = self.doc.Pages.pages[-1]
        last_page._ocg_properties = props

    def finalize_catalog(self):
        """在 canvas.save() 之前调用：将 OCProperties 写入 catalog。"""
        from reportlab.pdfbase import pdfdoc

        all_refs = [layer.ref for layer in self.layers]
        on_refs  = [layer.ref for layer in self.layers if layer.visible]

        # 图层显示顺序（必须是 OCG 间接引用，不能是字符串）
        order_items = [layer.ref for layer in self.layers]

        default_cfg = pdfdoc.PDFDictionary({
            "Name":      pdfdoc.PDFString("Default"),
            "BaseState": pdfdoc.PDFName("OFF"),
            "ON":        pdfdoc.PDFArray(on_refs),
            "Order":     pdfdoc.PDFArray(order_items),
        })

        oc_props = pdfdoc.PDFDictionary({
            "OCGs": pdfdoc.PDFArray(all_refs),
            "D":    default_cfg,
        })

        self.doc.Catalog.OCProperties = oc_props


# ─────────────────────────── PDF 构建 ───────────────────────────

# 目标页面宽度（points）：使用 A4 宽度 595pt ≈ 21cm
_TARGET_PAGE_WIDTH_PT = 595.0


def _compute_page_scale(pw_px: int, ph_px: int) -> tuple[float, float, float, float]:
    """
    将像素坐标映射到 PDF points。
    保持原始宽高比，目标宽度固定为 A4 宽度。
    返回 (pw_pt, ph_pt, sx, sy)
    """
    if pw_px <= 0 or ph_px <= 0:
        return _TARGET_PAGE_WIDTH_PT, 841.89, 1.0, 1.0
    sx = _TARGET_PAGE_WIDTH_PT / pw_px
    sy = sx  # 等比缩放
    return (pw_px * sx, ph_px * sy, sx, sy)


def draw_colored_text(c, x: float, y: float, text: str,
                      font_name: str, font_size: float,
                      default_color: tuple) -> None:
    """
    渲染可能含有 <red>...</red> / <blue>...</blue> 标记的文字。
    不含标记时按 default_color 绘制。
    """
    from reportlab.pdfbase import pdfmetrics

    # 解析内联颜色标记
    pattern = re.compile(r"<(red|blue)>(.*?)</\1>", re.DOTALL)
    parts = []
    last = 0
    for m in pattern.finditer(text):
        if m.start() > last:
            parts.append(("default", text[last:m.start()]))
        parts.append((m.group(1), m.group(2)))
        last = m.end()
    if last < len(text):
        parts.append(("default", text[last:]))

    if not parts:
        parts = [("default", text)]

    cx = x
    for color_key, chunk in parts:
        if not chunk:
            continue
        if color_key == "default":
            c.setFillColorRGB(*default_color)
        else:
            c.setFillColorRGB(*get_color_rgb(color_key))
        c.setFont(font_name, font_size)
        c.drawString(cx, y, chunk)
        cx += pdfmetrics.stringWidth(chunk, font_name, font_size)


def _wrap_text_lines(text: str, font_name: str, font_size: float,
                     max_width: float) -> list[str]:
    """将文字按最大宽度换行（支持中文按字符换行）。"""
    from reportlab.pdfbase import pdfmetrics

    # 先按已有换行符分段
    raw_lines = text.split("\n")
    result = []
    for raw_line in raw_lines:
        if not raw_line.strip():
            result.append("")
            continue
        # 去除内联标记来测量宽度
        clean = re.sub(r"<[^>]+>", "", raw_line)
        if pdfmetrics.stringWidth(clean, font_name, font_size) <= max_width:
            result.append(raw_line)
            continue
        # 逐字符换行
        current = ""
        current_clean = ""
        for char in raw_line:
            test_clean = current_clean + char
            if pdfmetrics.stringWidth(test_clean, font_name, font_size) <= max_width:
                current += char
                current_clean = test_clean
            else:
                if current:
                    result.append(current)
                current = char
                current_clean = char
        if current:
            result.append(current)
    return result


def build_pdf(pages: list[dict], output_path: Path, font_map: dict,
              line_style: str = "auto") -> None:
    """生成 5层可切换 PDF。"""
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.lib.utils import ImageReader

    _setup_ocg_support()

    print(f"[build] 生成 PDF → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    c = rl_canvas.Canvas(str(output_path))
    doc = c._doc
    ocg = OCGManager(c, doc)

    # 创建5个图层（resource_name, display_name, visible_by_default）
    L_SCAN   = ocg.add_layer("L1Scan",   "原始扫描件",     visible=False)
    L_ILLUS  = ocg.add_layer("L2Illus",  "增强插图",       visible=True)
    L_DIM    = ocg.add_layer("L3Dim",    "尺寸标注",       visible=True)
    L_TEXT   = ocg.add_layer("L4Text",   "数字化文字",     visible=True)
    L_LABEL  = ocg.add_layer("L5Labels", "标注引线",       visible=True)

    total_pages = len(pages)
    for idx, page in enumerate(pages):
        pnum = page["page_number"]
        pw_px = page.get("page_width_px",  908)
        ph_px = page.get("page_height_px", 1282)

        # 动态计算页面大小（保持宽高比，目标宽度 A4 = 595pt）
        pw_pt, ph_pt, sx, sy = _compute_page_scale(pw_px, ph_px)
        c.setPageSize((pw_pt, ph_pt))

        def px2pdf(bx, by, bw=None, bh=None):
            """像素坐标 → PDF坐标（左下角原点，Y轴翻转）"""
            px = bx * sx
            py = ph_pt - by * sy   # top-left pixel → bottom-left PDF
            if bw is not None and bh is not None:
                return (px, py - bh * sy, bw * sx, bh * sy)
            return (px, py)

        # ── Layer 0: 白色背景（始终可见） ──
        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, pw_pt, ph_pt, fill=1, stroke=0)

        # ── Layer 1: 原始扫描件 ──
        img_path = page.get("original_image_path")
        if img_path and os.path.exists(img_path):
            ocg.begin(L_SCAN.resource_name)
            try:
                c.drawImage(ImageReader(img_path), 0, 0,
                            width=pw_pt, height=ph_pt,
                            preserveAspectRatio=False)
            except Exception as e:
                print(f"  [警告] page_{pnum:03d}: 扫描件图片读取失败: {e}")
            ocg.end()

        # ── Layer 2: 增强插图 ──
        has_illus = False
        for illus in page.get("illustration_regions", []):
            ill_path = get_illustration_path(pnum, illus["id"], line_style)
            if ill_path:
                has_illus = True
                bx, by, bw, bh = illus["bbox"]
                ix, iy, iw, ih = px2pdf(bx, by, bw, bh)
                ocg.begin(L_ILLUS.resource_name)
                try:
                    c.drawImage(ImageReader(ill_path), ix, iy,
                                width=iw, height=ih,
                                preserveAspectRatio=True, mask="auto")
                except Exception as e:
                    print(f"  [警告] page_{pnum:03d}: 插图读取失败: {e}")
                ocg.end()

        # 按 font_level 分类文字区域
        dim_regions = [r for r in page["text_regions"] if r["font_level"] == "dimension"]
        label_regions = [r for r in page["text_regions"] if r["font_level"] == "label"]
        text_regions = [r for r in page["text_regions"]
                        if r["font_level"] not in ("dimension", "label")]

        # ── Layer 3: 尺寸标注 ──
        if dim_regions:
            ocg.begin(L_DIM.resource_name)
            for region in dim_regions:
                _draw_text_region(c, region, px2pdf, font_map, pw_pt, ph_pt, sx, sy)
            ocg.end()

        # ── Layer 4: 正文文字 ──
        if text_regions:
            ocg.begin(L_TEXT.resource_name)
            for region in text_regions:
                _draw_text_region(c, region, px2pdf, font_map, pw_pt, ph_pt, sx, sy)
            ocg.end()

        # ── Layer 5: 标注 ──
        if label_regions:
            ocg.begin(L_LABEL.resource_name)
            for region in label_regions:
                _draw_text_region(c, region, px2pdf, font_map, pw_pt, ph_pt, sx, sy)
            ocg.end()

        c.showPage()
        ocg.attach_to_page()  # 将 OCG 属性绑定到刚添加的页面

        if (idx + 1) % 10 == 0 or (idx + 1) == total_pages:
            print(f"  [PDF] {idx+1}/{total_pages} 页完成")

    # 写入 OCProperties 到 catalog，然后保存
    ocg.finalize_catalog()
    c.save()

    size_mb = output_path.stat().st_size / 1024 / 1024
    print(f"  [PDF] ✓ 保存: {output_path} ({size_mb:.1f} MB)")


def _draw_text_region(c, region: dict, px2pdf, font_map: dict,
                      pw_pt: float, ph_pt: float,
                      sx: float, sy: float) -> None:
    """在 PDF 画布上渲染一个文字区域（含换行）。"""
    from reportlab.pdfbase import pdfmetrics

    bx, by, bw, bh = region["bbox"]
    # PDF 坐标（左下角原点）
    rx = bx * sx
    ry = ph_pt - by * sy           # bbox 顶部对应的 PDF y
    rw = bw * sx
    rh = bh * sy

    # 每行使用第一行的 font_level 和 color
    lines_data = region.get("lines", [])
    if not lines_data:
        return

    for line_info in lines_data:
        content = line_info.get("content", "").strip()
        if not content:
            continue

        fl = line_info.get("font_level", region.get("font_level", "body"))
        color_name = line_info.get("color", region.get("color", "black"))

        font_name, base_size = font_map.get(fl, font_map["body"])
        color_rgb = get_color_rgb(color_name)

        # 字体大小：使用配置中的基础大小
        # 当 bbox 非常小时（< 1.5 行高），适当缩小以免溢出
        font_size = float(base_size)
        if rh > 0 and rh < font_size * 1.5:
            font_size = max(7.0, rh / 1.5)
        line_h = font_size * 1.4

        wrapped = _wrap_text_lines(content, font_name, font_size, rw)

        for li, wline in enumerate(wrapped):
            if not wline.strip():
                continue
            ly = ry - (li + 1) * line_h
            if ly < (ry - rh - font_size):   # 超出 bbox 不画
                break
            draw_colored_text(c, rx, ly, wline, font_name, font_size, color_rgb)


# ─────────────────���───────── PPTX 构建 ───────────────────────────

def build_pptx(pages: list[dict], output_path: Path, font_map: dict,
               line_style: str = "auto") -> None:
    """生成可编辑 PPTX。"""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    print(f"[build] 生成 PPTX → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # PPTX 字体名映射（系统/Noto）
    pptx_font_map = {
        "NotoSansSC-Bold":     "Noto Sans SC",
        "NotoSansSC-Regular":  "Noto Sans SC",
        "NotoSerifSC-Regular": "Noto Serif SC",
        "NotoSerifSC-Bold":    "Noto Serif SC",
        "STHeitiMedium":       "Heiti SC",
        "STHeitiLight":        "Heiti SC",
        "Songti":              "Songti SC",
        "Helvetica":           "Arial",
    }

    total_pages = len(pages)
    for idx, page in enumerate(pages):
        pnum = page["page_number"]
        pw_px = page.get("page_width_px", 908)
        ph_px = page.get("page_height_px", 1282)

        # 设置幻灯片尺寸（匹配页面比例，最大 13.33 英寸宽）
        max_w_emu = int(13.33 * 914400)
        max_h_emu = int(10.0  * 914400)
        aspect = pw_px / ph_px

        if aspect >= 1:
            slide_w = max_w_emu
            slide_h = int(slide_w / aspect)
        else:
            slide_h = max_h_emu
            slide_w = int(slide_h * aspect)

        prs.slide_width  = slide_w
        prs.slide_height = slide_h

        # 坐标转换（像素 → EMU）
        def px2emu(bx, by, bw=None, bh=None):
            ex = int(bx / pw_px * slide_w)
            ey = int(by / ph_px * slide_h)
            if bw is not None and bh is not None:
                ew = int(bw / pw_px * slide_w)
                eh = int(bh / ph_px * slide_h)
                return (ex, ey, ew, eh)
            return (ex, ey)

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 插图
        for illus in page.get("illustration_regions", []):
            ill_path = get_illustration_path(pnum, illus["id"], line_style)
            if ill_path:
                bx, by, bw, bh = illus["bbox"]
                ex, ey, ew, eh = px2emu(bx, by, bw, bh)
                try:
                    slide.shapes.add_picture(ill_path, left=ex, top=ey,
                                             width=ew, height=eh)
                except Exception as e:
                    print(f"  [警告] page_{pnum:03d}: 插图 PPTX 插入失败: {e}")

        # 文字区域
        for region in page["text_regions"]:
            bx, by, bw, bh = region["bbox"]
            ex, ey, ew, eh = px2emu(bx, by, bw, bh)

            # 最小文字框尺寸
            ew = max(ew, int(0.5 * 914400))
            eh = max(eh, int(0.3 * 914400))

            try:
                txbox = slide.shapes.add_textbox(left=ex, top=ey,
                                                  width=ew, height=eh)
                tf = txbox.text_frame
                tf.word_wrap = True

                lines_data = region.get("lines", [])
                for i, line_info in enumerate(lines_data):
                    content = line_info.get("content", "").strip()
                    if not content:
                        continue

                    fl = line_info.get("font_level", region.get("font_level", "body"))
                    color_name = line_info.get("color", region.get("color", "black"))

                    rl_font_name, pt_size = font_map.get(fl, font_map["body"])
                    pptx_fname = pptx_font_map.get(rl_font_name, "Microsoft YaHei")

                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    run = para.add_run()
                    # 去除内联标记
                    clean_content = re.sub(r"<[^>]+>", "", content)
                    run.text = clean_content
                    run.font.name = pptx_fname
                    run.font.size = Pt(pt_size)

                    cr, cg, cb = get_color_rgb(color_name)
                    run.font.color.rgb = RGBColor(
                        int(cr * 255), int(cg * 255), int(cb * 255)
                    )
                    if fl in ("title", "subtitle"):
                        run.font.bold = True

            except Exception as e:
                print(f"  [警告] page_{pnum:03d}: 文字框 {region['id']} 失败: {e}")

        # 备注
        conf = page.get("average_confidence")
        conf_str = f"{conf:.2f}" if conf else "N/A"
        illus_count = len(page.get("illustration_regions", []))
        notes = slide.notes_slide
        notes.notes_text_frame.text = (
            f"Page {pnum} | Confidence: {conf_str} | "
            f"Illustrations: {illus_count} | "
            f"Notes: {page.get('notes', '')[:200]}"
        )

        if (idx + 1) % 10 == 0 or (idx + 1) == total_pages:
            print(f"  [PPTX] {idx+1}/{total_pages} 页完成")

    prs.save(str(output_path))
    size_mb = output_path.stat().st_size / 1024 / 1024
    print(f"  [PPTX] ✓ 保存: {output_path} ({size_mb:.1f} MB)")


# ─────────────────────────── Markdown 构建 ───────────────────────────

def build_markdown(pages: list[dict], output_path: Path) -> None:
    """生成纯文字 Markdown。"""
    print(f"[build] 生成 Markdown → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    lines_out = [
        "# 教案数字版 — 纯文字",
        "",
        "> 自动生成，请校对。⚠️[?X?]⚠️ 标记为不确定字符。",
        "",
    ]

    for page in pages:
        pnum = page["page_number"]
        lines_out.append(f"---")
        lines_out.append(f"")
        lines_out.append(f"## 第 {pnum} 页")
        lines_out.append(f"")

        conf = page.get("average_confidence")
        if conf:
            conf_pct = int(conf * 100)
            lines_out.append(f"> 置信度: {conf_pct}%")
            lines_out.append(f"")

        # 插图占位
        for illus in page.get("illustration_regions", []):
            lines_out.append(
                f"[插图: {illus.get('illustration_type', '图示')} — "
                f"位置: {illus['bbox']}]"
            )
            lines_out.append("")

        # 文字区域（按 y 坐标排序）
        text_regions = sorted(page["text_regions"], key=lambda r: r["bbox"][1])
        for region in text_regions:
            content = region.get("content", "").strip()
            if not content:
                continue

            fl = region.get("font_level", "body")
            color = region.get("color", "black")

            # 处理内联颜色标记
            content = _md_colorize(content, color)

            # 替换不确定字符标记
            content = content.replace("[?", "⚠️[?").replace("?]", "?]⚠️")

            # 根据 font_level 格式化
            if fl == "title":
                lines_out.append(f"### {content}")
            elif fl == "subtitle":
                lines_out.append(f"#### {content}")
            elif fl == "dimension":
                lines_out.append(f"`{content}`")
            elif fl == "label":
                lines_out.append(f"**{content}**（标注）")
            elif fl == "annotation":
                lines_out.append(f"*{content}*")
            else:
                # body — 支持多行
                for ln in content.split("\n"):
                    if ln.strip():
                        lines_out.append(ln.strip())
            lines_out.append("")

    output_path.write_text("\n".join(lines_out), encoding="utf-8")
    size_kb = output_path.stat().st_size / 1024
    print(f"  [MD] ✓ 保存: {output_path} ({size_kb:.1f} KB)")


def _md_colorize(text: str, region_color: str) -> str:
    """
    将 <red>...</red> 标记转成 Markdown 粗体+注释，
    整体颜色由 region_color 决定。
    """
    # 处理内联标记
    text = re.sub(r"<red>(.*?)</red>",   r"**\1**（红色标注）", text, flags=re.DOTALL)
    text = re.sub(r"<blue>(.*?)</blue>", r"**\1**（蓝色标注）", text, flags=re.DOTALL)

    # 整体颜色注释
    if region_color == "red" and "<red>" not in text:
        text = f"**{text}**（红色）"
    elif region_color == "blue" and "<blue>" not in text:
        text = f"**{text}**（蓝色）"

    return text


# ─────────────────────────── QA 报告 ───────────────────────────

def build_qa_report(pages: list[dict], output_path: Path) -> None:
    """生成校对报告 Markdown。"""
    print(f"[build] 生成 QA 报告 → {output_path.name}")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    total_pages = len(pages)
    all_confs = [p["average_confidence"] for p in pages if p.get("average_confidence")]
    avg_conf = _mean(all_confs)

    # 统计不确定字符
    total_uncertain = 0
    total_japanese = 0
    low_conf_pages = []

    for page in pages:
        pnum = page["page_number"]
        conf = page.get("average_confidence", 0.0)
        if conf and conf < 0.75:
            low_conf_pages.append((pnum, conf))

        for region in page["text_regions"]:
            content = region.get("content", "")
            total_uncertain += content.count("[?")
            total_japanese += len(re.findall(r"[\u3040-\u309F\u30A0-\u30FF]", content))

    lines = [
        "# 校对报告",
        "",
        "## 汇总统计",
        "",
        f"| 项目 | 数值 |",
        f"|------|------|",
        f"| 总页数 | {total_pages} |",
        f"| 平均置信度 | {f'{avg_conf:.1%}' if avg_conf else 'N/A'} |",
        f"| 不确定字符数 | {total_uncertain} |",
        f"| 检测到日文字符数 | {total_japanese} |",
        f"| 低置信度页面（<75%） | {len(low_conf_pages)} |",
        "",
    ]

    if low_conf_pages:
        lines += [
            "## 低置信度页面（需优先校对）",
            "",
            "| 页码 | 置信度 |",
            "|------|--------|",
        ]
        for pnum, conf in sorted(low_conf_pages, key=lambda x: x[1]):
            lines.append(f"| 第 {pnum} 页 | {conf:.1%} |")
        lines.append("")

    lines += [
        "## 各页详情",
        "",
        "| 页码 | 置信度 | 文字区域数 | 插图数 | 不确定字符 |",
        "|------|--------|-----------|--------|-----------|",
    ]

    for page in pages:
        pnum = page["page_number"]
        conf = page.get("average_confidence")
        conf_str = f"{conf:.1%}" if conf else "N/A"
        text_count = len(page["text_regions"])
        illus_count = len(page.get("illustration_regions", []))
        uncertain = sum(r.get("content", "").count("[?") for r in page["text_regions"])
        lines.append(f"| 第 {pnum} 页 | {conf_str} | {text_count} | {illus_count} | {uncertain} |")

    lines += ["", "---", "", "_报告由 build.py 自动生成_", ""]

    output_path.write_text("\n".join(lines), encoding="utf-8")
    size_kb = output_path.stat().st_size / 1024
    print(f"  [QA] ✓ 保存: {output_path} ({size_kb:.1f} KB)")


# ─────────────────────────── 页码范围解析 ───────────────────────────

def parse_pages_arg(arg: str) -> list[int]:
    pages = set()
    for part in arg.split(","):
        part = part.strip()
        if "-" in part:
            lo, hi = part.split("-", 1)
            pages.update(range(int(lo), int(hi) + 1))
        else:
            pages.add(int(part))
    return sorted(pages)


# ─────────────────────────── 主流程 ───────────────────��───────

def main():
    parser = argparse.ArgumentParser(description="Phase 3: 文档合成")
    parser.add_argument("--pages", type=str, default=None,
                        help='指定页码范围，如 "1-5" / "2,4,6"')
    parser.add_argument("--format", type=str, default="pdf,pptx,text,qa",
                        help="输出格式，逗号分隔（pdf,pptx,text,qa）")
    parser.add_argument("--line-style",
                        choices=["auto", "original", "strengthen", "vectorize", "inpainted"],
                        default="auto",
                        help="插图线稿样式（默认 auto）")
    parser.add_argument("--force", action="store_true",
                        help="强制覆盖已有输出文件")
    parser.add_argument("--download-fonts", action="store_true",
                        help="下载 Noto CJK 字体（约 50MB/个）")
    args = parser.parse_args()

    formats = {f.strip() for f in args.format.split(",")}
    page_nums = parse_pages_arg(args.pages) if args.pages else None

    # 加载数据
    pages = load_page_data(page_nums)
    if not pages:
        print("[错误] 没有找到任何页面数据，请先运行前序步骤。")
        sys.exit(1)

    # 初始化字体
    font_map = setup_fonts(download_noto=args.download_fonts)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    # 输出路径
    pdf_out  = OUTPUT_DIR / "教案_数字版.pdf"
    pptx_out = OUTPUT_DIR / "教案_数字版.pptx"
    md_out   = OUTPUT_DIR / "教案_纯文字.md"
    qa_out   = OUTPUT_DIR / "校对报告.md"

    # 生成各输出
    if "pdf" in formats:
        if not pdf_out.exists() or args.force:
            build_pdf(pages, pdf_out, font_map, line_style=args.line_style)
        else:
            print(f"[build] 跳过 PDF（已存在，使用 --force 覆盖）")

    if "pptx" in formats:
        if not pptx_out.exists() or args.force:
            build_pptx(pages, pptx_out, font_map, line_style=args.line_style)
        else:
            print(f"[build] 跳过 PPTX（已存在，使用 --force 覆盖）")

    if "text" in formats:
        if not md_out.exists() or args.force:
            build_markdown(pages, md_out)
        else:
            print(f"[build] 跳过 Markdown（已存在，使用 --force 覆盖）")

    if "qa" in formats:
        if not qa_out.exists() or args.force:
            build_qa_report(pages, qa_out)
        else:
            print(f"[build] 跳过 QA 报告（已存在，使用 --force 覆盖）")

    print()
    print("─" * 50)
    print("[build] 完成！输出目录:", OUTPUT_DIR)
    for p in [pdf_out, pptx_out, md_out, qa_out]:
        if p.exists():
            size = p.stat().st_size
            unit = "MB" if size > 1024*1024 else "KB"
            val  = size/1024/1024 if size > 1024*1024 else size/1024
            print(f"  {p.name}: {val:.1f} {unit}")
    print("─" * 50)


if __name__ == "__main__":
    main()
